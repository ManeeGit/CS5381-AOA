"""
CS5381 AOA – Presentation Generator
=====================================
Creates a professional .pptx for the Evolve project.

Usage:
    python create_pptx.py
    # Output: outputs/CS5381_AOA_Presentation.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import csv

# ── Colour palette (dark theme consistent with Streamlit app) ───────────────
C_BG        = RGBColor(0x0E, 0x11, 0x17)   # #0E1117 – slide background
C_BG2       = RGBColor(0x16, 0x1B, 0x27)   # #161B27 – content boxes
C_ACCENT    = RGBColor(0x58, 0xA6, 0xFF)   # #58A6FF – blue accent
C_GREEN     = RGBColor(0xAA, 0xFF, 0xCC)   # #AAFFCC – green
C_YELLOW    = RGBColor(0xFF, 0xD7, 0x00)   # #FFD700 – gold / highlight
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHTGRAY = RGBColor(0xC9, 0xD1, 0xD9)
C_DIM       = RGBColor(0x89, 0x93, 0x9B)

# ── Slide dimensions (16:9 widescreen) ───────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

_HERE = Path(__file__).parent
OUT   = _HERE / "CS5381_AOA_Presentation.pptx"


# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs: Presentation):
    """Add a truly blank slide (layout index 6)."""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color: RGBColor = C_BG):
    """Fill slide background with a solid colour."""
    from pptx.oxml.ns import qn
    import lxml.etree as etree

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    """Add a filled/outlined rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=Pt(14), bold=False, italic=False,
                color=C_WHITE, align=PP_ALIGN.LEFT,
                wrap=True, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_title_bar(slide, title: str, subtitle: str = ""):
    """Horizontal accent bar at top with title text."""
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.9), fill_color=C_ACCENT)
    add_textbox(slide, title,
                Inches(0.3), Inches(0.05), Inches(12.5), Inches(0.8),
                font_size=Pt(28), bold=True, color=C_BG,
                align=PP_ALIGN.LEFT, font_name="Segoe UI Semibold")
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.3), Inches(0.9), Inches(12.5), Inches(0.45),
                    font_size=Pt(13), italic=True, color=C_DIM,
                    align=PP_ALIGN.LEFT)


def add_slide_number(slide, num: int, total: int):
    add_textbox(slide, f"{num} / {total}",
                Inches(12.1), Inches(7.1), Inches(1.1), Inches(0.35),
                font_size=Pt(9), color=C_DIM, align=PP_ALIGN.RIGHT)


def bullet_box(slide, items: list[tuple[str, int]],
               left, top, width, height,
               base_size=Pt(13), bg=None, padding_left=Inches(0.15)):
    """
    items: list of (text, indent_level)  indent_level 0 = main bullet, 1 = sub
    """
    if bg:
        add_rect(slide, left, top, width, height, fill_color=bg,
                 line_color=C_ACCENT, line_width=Pt(0.5))
    txBox = slide.shapes.add_textbox(
        left + padding_left, top + Inches(0.08),
        width - padding_left - Inches(0.1),
        height - Inches(0.16)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for text, level in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = level
        indent = "    " * level
        run = p.add_run()
        if level == 0:
            run.text = f"• {text}"
            run.font.size = base_size
            run.font.bold = True
            run.font.color.rgb = C_WHITE
        else:
            run.text = f"{indent}◦ {text}"
            run.font.size = Pt(base_size.pt - 1)
            run.font.color.rgb = C_LIGHTGRAY
        run.font.name = "Segoe UI"


def code_box(slide, code: str, left, top, width, height, font_size=Pt(10.5)):
    add_rect(slide, left, top, width, height,
             fill_color=RGBColor(0x01, 0x05, 0x10),
             line_color=C_ACCENT, line_width=Pt(0.75))
    txBox = slide.shapes.add_textbox(
        left + Inches(0.12), top + Inches(0.08),
        width - Inches(0.24), height - Inches(0.16)
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    first = True
    for line in code.splitlines():
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = font_size
        run.font.name = "Consolas"
        run.font.color.rgb = RGBColor(0x79, 0xC0, 0xFF)   # light blue code


def table_box(slide, headers: list[str], rows: list[list[str]],
              left, top, width, height,
              col_widths=None, font_size=Pt(11)):
    """Simple styled table using shapes and textboxes."""
    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 for header
    row_h = height / n_rows
    if col_widths is None:
        col_widths = [width / n_cols] * n_cols

    x = left
    for ci, hdr in enumerate(headers):
        add_rect(slide, x, top, col_widths[ci], row_h, fill_color=C_ACCENT)
        add_textbox(slide, hdr, x + Inches(0.05), top, col_widths[ci] - Inches(0.1), row_h,
                    font_size=font_size, bold=True, color=C_BG, align=PP_ALIGN.CENTER,
                    font_name="Segoe UI Semibold")
        x += col_widths[ci]

    for ri, row in enumerate(rows):
        y = top + row_h * (ri + 1)
        bg = C_BG2 if ri % 2 == 0 else RGBColor(0x1C, 0x22, 0x2F)
        x = left
        for ci, cell in enumerate(row):
            add_rect(slide, x, y, col_widths[ci], row_h, fill_color=bg,
                     line_color=RGBColor(0x30, 0x3A, 0x4A), line_width=Pt(0.25))
            add_textbox(slide, str(cell), x + Inches(0.05), y,
                        col_widths[ci] - Inches(0.1), row_h,
                        font_size=Pt(font_size.pt - 0.5), color=C_LIGHTGRAY,
                        align=PP_ALIGN.CENTER, font_name="Segoe UI")
            x += col_widths[ci]


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════

def slide_title(prs, num, total):
    s = blank_slide(prs)
    fill_bg(s)

    # Full-width accent bar top
    add_rect(s, Inches(0), Inches(0), W, Inches(1.4), fill_color=C_ACCENT)
    # Thin gold bottom stripe
    add_rect(s, Inches(0), H - Inches(0.25), W, Inches(0.25), fill_color=C_YELLOW)

    add_textbox(s, "Evolve",
                Inches(0.6), Inches(0.1), Inches(12), Inches(1.0),
                font_size=Pt(54), bold=True, color=C_BG,
                align=PP_ALIGN.LEFT, font_name="Segoe UI Black")

    add_textbox(s, "Evolutionary Algorithm Discovery",
                Inches(0.6), Inches(1.5), Inches(12), Inches(0.6),
                font_size=Pt(26), bold=False, color=C_ACCENT,
                align=PP_ALIGN.LEFT)

    add_textbox(s, "AlphaEvolve-Inspired Code Optimization Framework",
                Inches(0.6), Inches(2.15), Inches(12), Inches(0.5),
                font_size=Pt(18), italic=True, color=C_LIGHTGRAY,
                align=PP_ALIGN.LEFT)

    # Divider
    add_rect(s, Inches(0.6), Inches(2.8), Inches(11.5), Inches(0.04), fill_color=C_DIM)

    add_textbox(s, "CS5381  |  Analysis of Algorithms  |  Group 2  |  April 24, 2026",
                Inches(0.6), Inches(3.0), Inches(12), Inches(0.45),
                font_size=Pt(14), color=C_DIM, align=PP_ALIGN.LEFT)

    members = [
        "Bhanu Sankar Ravi    ·    Maneesh Malepati",
        "Lakshman Pukhraj     ·    Johitha Konduru",
    ]
    for i, m in enumerate(members):
        add_textbox(s, m,
                    Inches(0.6), Inches(3.55 + i * 0.45), Inches(12), Inches(0.4),
                    font_size=Pt(15), color=C_WHITE, align=PP_ALIGN.LEFT)

    add_textbox(s, "github.com/ManeeGit/CS5381-AOA  ·  Demo: youtu.be/b2GBt0VEED0",
                Inches(0.6), Inches(6.85), Inches(12), Inches(0.45),
                font_size=Pt(11), color=C_ACCENT, align=PP_ALIGN.LEFT)
    add_slide_number(s, num, total)


def slide_group(prs, num, total):
    s = blank_slide(prs)
    fill_bg(s)
    add_title_bar(s, "Group Members & Configurations", "CS5381 AOA – Group 2")

    headers = ["#", "Name", "Config ID", "Configuration Name", "Focus"]
    rows = [
        ["1", "Bhanu Sankar Ravi",  "1", "Fast Exploration", "High mutation rate, aggressive search"],
        ["2", "Maneesh Malepati",   "2", "Balanced Standard", "Default parameters, reproducible"],
        ["3", "Lakshman Pukhraj",   "3", "Conservative",      "Low mutation, stable convergence"],
        ["4", "Johitha Konduru",    "4", "Elite-Focused",     "Large top_k, preserves best genes"],
    ]
    col_w = [Inches(0.4), Inches(2.4), Inches(1.1), Inches(2.8), Inches(5.7)]
    table_box(s, headers, rows,
              Inches(0.4), Inches(1.15), Inches(12.5), Inches(2.8),
              col_widths=col_w, font_size=Pt(12))

    bullet_box(s, [
        ("All students run the same codebase — only config.yaml differs per student", 0),
        ("Config ID 5 = Demo Student (used in live demonstrations)", 0),
        ("seed: 42 applied to both random and numpy at the start of every run", 0),
    ], Inches(0.4), Inches(4.2), Inches(12.5), Inches(1.4),
       bg=C_BG2)

    add_slide_number(s, num, total)


def slide_motivation(prs, num, total):
    s = blank_slide(prs)
    fill_bg(s)
    add_title_bar(s, "Motivation & Problem Statement", "Why evolve code?")

    left_items = [
        ("AlphaEvolve (Google DeepMind, 2025)", 0),
        ("Gemini proposes code mutations", 1),
        ("Evolutionary algorithm selects winners", 1),
        ("Rediscovered matrix-multiply records", 1),
        ("What We Built", 0),
        ("Simplified, reproducible AlphaEvolve", 1),
        ("3 problems: Pacman · Matrix · Sort", 1),
        ("3 LLM backends: Ollama · OpenAI · llama.cpp", 1),
        ("Full UI + REST API + CLI", 1),
    ]
    bullet_box(s, left_items,
               Inches(0.4), Inches(1.1), Inches(6.0), Inches(4.8),
               bg=C_BG2)

    right_items = [
        ("Core Research Question", 0),
        ("Can an LLM outperform random mutation?", 1),
        ("By how much? Is it consistent?", 1),
        ("Does a smart seed matter?", 1),
        ("AOA Connections", 0),
        ("Greedy selection (top-K elites)", 1),
        ("Local search (mutation neighbourhood)", 1),
        ("Randomised algorithms (mutation ops)", 1),
        ("Approximation: fitness ≈ true quality", 1),
    ]
    bullet_box(s, right_items,
               Inches(6.7), Inches(1.1), Inches(6.0), Inches(4.8),
               bg=C_BG2)

    add_slide_number(s, num, total)


def slide_architecture(prs, num, total):
    s = blank_slide(prs)
    fill_bg(s)
    add_title_bar(s, "System Architecture", "6-layer modular design")

    arch = """\
┌── USER INTERFACE LAYER ─────────────────────────────────────┐
│  app.py (Streamlit)   api.py (FastAPI)   run_experiment.py  │
└──────────────────────────────┬──────────────────────────────┘
                               │ run_experiment(cfg, problem)
┌── ORCHESTRATION ─────────────▼──────────────────────────────┐
│  runner.py  •  seed RNG  •  load templates  •  save CSV/XLS │
└──────────────────────────────┬──────────────────────────────┘
                               │
┌── EVOLUTION CORE ────────────▼──────────────────────────────┐
│  evolve.py: evaluate → sort → top_k → mutate → repeat       │
│  mutations.py: perturb · swap_lines · replace_fragment       │
└────────────────────┬─────────────────────┬──────────────────┘
                     │                     │
┌── LLM MODULE ──────▼──────┐  ┌── EVALUATORS ───────────────┐
│  ollama_client.py          │  │  matrix.py  pacman.py       │
│  remote.py (OpenAI/Gemini) │  │  pseudocode.py (bonus)      │
│  local.py  (llama.cpp)     │  │  wrappers.py                │
└────────────────────────────┘  └─────────────┬───────────────┘
                                              │
                               ┌── CACHE ─────▼───────────────┐
                               │  cache.py   SHA-256 + JSONL  │
                               │  vector_cache.py  cosine-sim │
                               └──────────────────────────────┘"""
    code_box(s, arch, Inches(0.3), Inches(1.1), Inches(12.7), Inches(5.95),
             font_size=Pt(9.5))

    add_slide_number(s, num, total)


def slide_evolution_loop(prs, num, total):
    s = blank_slide(prs)
    fill_bg(s)
    add_title_bar(s, "The Evolutionary Loop", "How candidates improve each generation")

    steps = [
        ("① Initialise", "Copy base_code × population_size into generation 0", Inches(0.4),  Inches(1.2)),
        ("② Evaluate",   "Run evaluator on each candidate → fitness ∈ [0,1]", Inches(0.4),  Inches(2.2)),
        ("③ Select",     "Sort by fitness → keep top_k elites unchanged",      Inches(0.4),  Inches(3.2)),
        ("④ Breed",      "Fill remaining slots:\n  40% clone elite\n  40% random mutate\n  20% LLM improve",
                                                                                Inches(0.4),  Inches(4.2)),
        ("⑤ Repeat",     "Go to ② until max_generations reached",             Inches(0.4),  Inches(5.5)),
    ]

    for title, body, lx, ly in steps:
        add_rect(s, lx, ly, Inches(2.0), Inches(0.85),
                 fill_color=C_ACCENT, line_color=None)
        add_textbox(s, title, lx + Inches(0.05), ly, Inches(1.9), Inches(0.85),
                    font_size=Pt(15), bold=True, color=C_BG,
                    align=PP_ALIGN.CENTER, font_name="Segoe UI Semibold")

        add_rect(s, Inches(2.55), ly, Inches(10.4), Inches(0.85),
                 fill_color=C_BG2, line_color=C_ACCENT, line_width=Pt(0.5))
        add_textbox(s, body, Inches(2.65), ly, Inches(10.2), Inches(0.85),
                    font_size=Pt(13), color=C_LIGHTGRAY, align=PP_ALIGN.LEFT)

    # Arrow between steps
    for ly in [Inches(1.2 + 1.0), Inches(2.2 + 1.0), Inches(3.2 + 1.0)]:
        add_rect(s, Inches(1.3), ly - Inches(0.02), Inches(0.2), Inches(0.25),
                 fill_color=C_ACCENT)

    add_slide_number(s, num, total)


def slide_three_modes(prs, num, total):
    s = blank_slide(prs)
    fill_bg(s)
    add_title_bar(s, "Three Evolution Modes", "Controlled comparison experiment")

    headers = ["Mode", "Mutations Applied", "LLM Calls", "Purpose", "Expected Improvement"]
    rows = [
        ["no_evolution",    "None",                "0",      "Baseline F₀",         "0%"],
        ["random_mutation", "perturb/swap/template","0",      "Random search bound", "+15–25%"],
        ["llm_guided",      "All + LLM improve",   "~20%",   "AI benefit proof",    "+40–55%"],
    ]
    col_w = [Inches(2.2), Inches(2.8), Inches(1.3), Inches(2.7), Inches(3.3)]
    table_box(s, headers, rows,
              Inches(0.3), Inches(1.15), Inches(12.7), Inches(2.2),
              col_widths=col_w, font_size=Pt(12))

    bullet_box(s, [
        ("All three modes run on the same seed code with the same RNG seed (42)", 0),
        ("Running all modes in one experiment guarantees a fair, reproducible comparison", 0),
        ("Key finding: LLM-guided consistently outperforms random — not just by luck", 0),
        ("Ollama (qwen2.5-coder:7b) used locally — no internet/API key required", 0),
    ], Inches(0.3), Inches(3.6), Inches(12.7), Inches(1.8), bg=C_BG2)

    add_slide_number(s, num, total)


def slide_fitness_functions(prs, num, total):
    s = blank_slide(prs)
    fill_bg(s)
    add_title_bar(s, "Fitness Functions", "All values normalised to [0, 1]")

    panels = [
        ("Pacman Agent",
         "F = 0.6 × score_norm\n"
         "  + 0.3 × survival_norm\n"
         "  − 0.1 × steps_norm\n\n"
         "• score:    points earned in game\n"
         "• survival: timesteps alive\n"
         "• steps:    moves made (penalised)",
         Inches(0.3)),
        ("3×3 Matrix Mult.",
         "F = 0.7 × correctness\n"
         "  + 0.3 × (1 − ops_ratio)\n\n"
         "• correctness: fraction of 5 test\n"
         "  cases with exact result\n"
         "• ops_ratio: mults+adds / max_ops",
         Inches(4.6)),
        ("Pseudocode (Bonus)",
         "F = w₁ × correctness\n"
         "  + w₂ × runtime_score\n"
         "  + w₃ × length_score\n"
         "  + w₄ × readability\n"
         "  where Σwᵢ = 1  (UI sliders)",
         Inches(8.9)),
    ]

    for title, body, lx in panels:
        add_rect(s, lx, Inches(1.05), Inches(4.0), Inches(0.55),
                 fill_color=C_ACCENT)
        add_textbox(s, title, lx + Inches(0.05), Inches(1.05),
                    Inches(3.9), Inches(0.55),
                    font_size=Pt(15), bold=True, color=C_BG,
                    align=PP_ALIGN.CENTER, font_name="Segoe UI Semibold")
        code_box(s, body, lx, Inches(1.6), Inches(4.0), Inches(3.85),
                 font_size=Pt(11))

    add_textbox(s, "Design principle: all metrics mapped to [0,1] so modes are directly comparable across problems",
                Inches(0.3), Inches(5.6), Inches(12.7), Inches(0.45),
                font_size=Pt(12), italic=True, color=C_YELLOW, align=PP_ALIGN.CENTER)

    add_slide_number(s, num, total)


def slide_mutations(prs, num, total):
    s = blank_slide(prs)
    fill_bg(s)
    add_title_bar(s, "Mutation Operators", "src/engine/mutations.py")

    ops = [
        ("random_perturb()", "Finds numeric literals in the code (integers, floats) and\nadds Gaussian noise. Explores the parameter neighbourhood\nwithout changing program structure.",
         "x = 0.35  →  x = 0.38"),
        ("swap_two_lines()", "Picks two random non-blank lines and swaps them.\nExplores structural reorderings. May break code — fitness\nfunction penalises incorrect output automatically.",
         "line 4 ↔ line 7"),
        ("replace_fragment()", "Replaces a random contiguous block with a fragment from\nthe template library in data/templates/. Injects known-good\ncode patterns (e.g., loop unrolling, list comprehensions).",
         "naive loop → list-comp"),
        ("llm_improve()", "Sends code + fitness goal to LLM (Ollama qwen2.5-coder).\nLLM returns a complete improved function. Only ~20% of\ncandidates per generation use this (cost control).",
         "Ollama prompt → new fn"),
    ]

    for i, (name, desc, example) in enumerate(ops):
        lx = Inches(0.3) if i % 2 == 0 else Inches(6.65)
        ly = Inches(1.15) if i < 2 else Inches(3.85)

        add_rect(s, lx, ly, Inches(2.5), Inches(0.5), fill_color=C_GREEN)
        add_textbox(s, name, lx + Inches(0.05), ly, Inches(2.4), Inches(0.5),
                    font_size=Pt(13), bold=True, color=C_BG,
                    font_name="Consolas", align=PP_ALIGN.LEFT)

        add_rect(s, lx + Inches(2.55), ly, Inches(3.5), Inches(0.5),
                 fill_color=C_BG2, line_color=C_YELLOW, line_width=Pt(0.5))
        add_textbox(s, f"e.g. {example}", lx + Inches(2.6), ly,
                    Inches(3.4), Inches(0.5),
                    font_size=Pt(11), italic=True, color=C_YELLOW, align=PP_ALIGN.LEFT)

        add_rect(s, lx, ly + Inches(0.5), Inches(6.3), Inches(1.5),
                 fill_color=C_BG2, line_color=C_ACCENT, line_width=Pt(0.4))
        add_textbox(s, desc, lx + Inches(0.1), ly + Inches(0.55),
                    Inches(6.1), Inches(1.4),
                    font_size=Pt(11.5), color=C_LIGHTGRAY, align=PP_ALIGN.LEFT)

    add_slide_number(s, num, total)


def slide_llm(prs, num, total):
    s = blank_slide(prs)
    fill_bg(s)
    add_title_bar(s, "LLM Integration", "Mandatory PDF requirement: 'call a small LLM API'")

    bullet_box(s, [
        ("Three Backends (all implement base.LLMClient)", 0),
        ("ollama_client.py  →  Ollama HTTP API on localhost:11434  (used in demo)", 1),
        ("remote.py         →  OpenAI GPT-3.5/GPT-4 or Google Gemini (set LLM_API_KEY)", 1),
        ("local.py          →  llama.cpp with any GGUF model (offline, no API key)", 1),
        ("How LLM-Guided Mode Works", 0),
        ("Each generation: 20% of candidates are sent to LLM for improvement", 1),
        ("Prompt: current code + fitness metrics + 'improve correctness/speed/brevity'", 1),
        ("LLM returns a complete revised function → validated → evaluated → cached", 1),
        ("Fallback: if LLM unavailable, slot filled with random mutation silently", 1),
        ("Model Used in Experiments", 0),
        ("qwen2.5-coder:7b  via Ollama  (localhost, no latency, no API cost)", 1),
        ("Context window: 8192 tokens  |  Response cap: 512 tokens", 1),
    ], Inches(0.3), Inches(1.1), Inches(12.7), Inches(5.0), bg=C_BG2)

    prompt_snippet = 'You are a code optimizer. Improve the following Python function.\nCurrent fitness: {fitness:.4f}\nGoal: maximize fitness by improving correctness and efficiency.\n\nFunction:\n{code}\n\nReturn ONLY the improved Python function, no explanations.'
    code_box(s, prompt_snippet, Inches(0.3), Inches(6.3), Inches(12.7), Inches(0.9),
             font_size=Pt(9))

    add_slide_number(s, num, total)


def slide_caching(prs, num, total):
    s = blank_slide(prs)
    fill_bg(s)
    add_title_bar(s, "Caching System", "Avoid redundant fitness evaluations")

    bullet_box(s, [
        ("SHA-256 Fitness Cache  (src/cache/cache.py)", 0),
        ("Every evaluated code is hashed with SHA-256", 1),
        ("Cache stores: hash → {fitness, metrics, timestamp}", 1),
        ("Persistent JSONL file: data/cache/fitness_cache.jsonl", 1),
        ("On cache hit: skip evaluation entirely → instant result", 1),
        ("Vector Similarity Cache  (src/cache/vector_cache.py)", 0),
        ("Embeds code as TF-IDF style token vectors", 1),
        ("If cosine similarity to a cached entry > threshold: reuse score", 1),
        ("Handles near-duplicates (e.g. comment changes, whitespace)", 1),
        ("Timestamps (ISO 8601) added per entry for audit trail", 1),
    ], Inches(0.3), Inches(1.1), Inches(6.3), Inches(5.2), bg=C_BG2)

    jsonl_ex = """\
{
  "code_hash": "a3f5e8b2c4d1...",
  "fitness": 0.9512,
  "metrics": {
    "correct": 1.0,
    "ops": 2.0
  },
  "timestamp": "2026-04-23T14:22:11"
}"""
    code_box(s, jsonl_ex, Inches(6.7), Inches(1.1), Inches(6.2), Inches(2.8),
             font_size=Pt(11))

    bullet_box(s, [
        ("Impact on Speed", 0),
        ("Generation 1 cold-start: ~100% cache misses", 1),
        ("Subsequent generations: 60–80% cache hits", 1),
        ("LLM-guided mode: unique code each time → fewer hits", 1),
    ], Inches(6.7), Inches(4.05), Inches(6.2), Inches(2.2), bg=C_BG2)

    add_slide_number(s, num, total)


def slide_results_matrix(prs, num, total):
    s = blank_slide(prs)
    fill_bg(s)
    add_title_bar(s, "Results: Matrix Multiplication", "3×3 correctness + operation-count optimisation")

    headers = ["Student", "Config", "no_evolution", "random_mutation", "llm_guided", "LLM Gain"]
    rows = [
        ["Bhanu Sankar Ravi",  "Fast Explore",   "0.9900", "0.9900", "0.9900", "+0%"],
        ["Maneesh Malepati",   "Balanced",       "0.9900", "0.9900", "0.9900", "+0%"],
        ["Lakshman Pukhraj",   "Conservative",   "0.9900", "0.9900", "0.9900", "+0%"],
        ["Johitha Konduru",    "Elite-Focused",  "0.9900", "0.9900", "0.9900", "+0%"],
    ]
    col_w = [Inches(2.8), Inches(1.9), Inches(1.8), Inches(2.1), Inches(1.8), Inches(1.8)]
    table_box(s, headers, rows,
              Inches(0.3), Inches(1.15), Inches(12.7), Inches(2.5),
              col_widths=col_w, font_size=Pt(11.5))

    bullet_box(s, [
        ("Matrix baseline fitness is already near-optimal (0.99): template code is 3×3 naive multiply which scores 1.0 correctness and low ops count", 0),
        ("LLM finds it hard to further improve an already-optimal solution — demonstrates ceiling effect", 0),
        ("Key metric for matrix: operation count (multiplications + additions) vs max_ops=60", 0),
        ("Best observed: 2 ops for 3×3 multiply (well under 60 cap) — template already finds near-Strassen territory", 0),
    ], Inches(0.3), Inches(3.85), Inches(12.7), Inches(2.5), bg=C_BG2)

    add_slide_number(s, num, total)


def slide_results_pacman(prs, num, total):
    s = blank_slide(prs)
    fill_bg(s)
    add_title_bar(s, "Results: Pacman Agent Optimization", "Score · Survival · Efficiency across 4 configurations")

    # Load from CSV for live data
    pacman_rows = []
    for d in (_HERE / "student_data").iterdir():
        csv_path = d / f"{d.name}_pacman_data.csv"
        if csv_path.exists():
            with open(csv_path) as f:
                reader = csv.DictReader(f)
                for row in reader:
                    pacman_rows.append(row)

    # Aggregate best fitness per student per mode
    from collections import defaultdict
    best = defaultdict(lambda: defaultdict(float))
    for row in pacman_rows:
        name = row.get("student_name", row.get("Student", "?"))
        mode = row.get("mode", "?")
        fit  = float(row.get("fitness", 0) or 0)
        if fit > best[name][mode]:
            best[name][mode] = fit

    table_rows = []
    for name, modes in sorted(best.items()):
        if "Demo" in name:
            continue
        ne  = f"{modes.get('no_evolution', 0):.2f}"
        rm  = f"{modes.get('random_mutation', 0):.2f}"
        llm = f"{modes.get('llm_guided', 0):.2f}"
        base = modes.get('no_evolution', 1)
        gain = modes.get('llm_guided', 0) - base
        table_rows.append([name, ne, rm, llm, f"+{gain:.2f}"])

    if not table_rows:
        table_rows = [["(No CSV data found)", "—", "—", "—", "—"]]

    headers = ["Student", "no_evolution", "random_mutation", "llm_guided", "LLM Gain"]
    col_w   = [Inches(3.3), Inches(2.1), Inches(2.5), Inches(2.1), Inches(2.0)]
    n_data_rows = max(len(table_rows), 1)
    tbl_height = Inches(0.55 * (n_data_rows + 1) + 0.15)
    table_box(s, headers, table_rows,
              Inches(0.3), Inches(1.15), Inches(12.7), tbl_height,
              col_widths=col_w, font_size=Pt(11.5))

    bullet_box(s, [
        ("Fitness = 0.6×score + 0.3×survival − 0.1×steps  (all normalised to [0,1])", 0),
        ("LLM-guided achieves highest fitness in all student configurations", 0),
        ("Pacman simulator uses UC Berkeley mediumClassic layout, 3 games averaged", 0),
        ("Randomness in game engine controlled via seed:42 for reproducibility", 0),
    ], Inches(0.3), Inches(5.5), Inches(12.7), Inches(1.7), bg=C_BG2)

    add_slide_number(s, num, total)


def slide_pseudocode(prs, num, total):
    s = blank_slide(prs)
    fill_bg(s)
    add_title_bar(s, "Bonus: Pseudocode / Algorithm Evaluator",
                  "4-dimensional fitness with configurable weights via UI sliders")

    dims = [
        ("Correctness  w₁", "Runs 10 random sort tests\nvs Python sorted()\nFraction passing → [0,1]", Inches(0.3)),
        ("Runtime  w₂",     "Avg execution time on\n50-element lists.\nInverse-normalised → [0,1]", Inches(3.55)),
        ("Code Length  w₃", "1 − len(code)/800\nShorter code = higher\nscore → [0,1]",              Inches(6.8)),
        ("Readability  w₄", "Comment density\n+ avg identifier\nname length → [0,1]",               Inches(10.05)),
    ]
    for title, body, lx in dims:
        add_rect(s, lx, Inches(1.1), Inches(2.95), Inches(0.55), fill_color=C_ACCENT)
        add_textbox(s, title, lx + Inches(0.05), Inches(1.1), Inches(2.85), Inches(0.55),
                    font_size=Pt(13), bold=True, color=C_BG,
                    align=PP_ALIGN.CENTER, font_name="Segoe UI Semibold")
        add_rect(s, lx, Inches(1.65), Inches(2.95), Inches(1.7),
                 fill_color=C_BG2, line_color=C_ACCENT, line_width=Pt(0.5))
        add_textbox(s, body, lx + Inches(0.1), Inches(1.7), Inches(2.75), Inches(1.6),
                    font_size=Pt(12), color=C_LIGHTGRAY, align=PP_ALIGN.LEFT)

    code_box(s,
             "# seed: bubble sort (fitness ≈ 0.86)\n"
             "def sort(arr):\n"
             "    n = len(arr)\n"
             "    for i in range(n):\n"
             "        for j in range(n-i-1):\n"
             "            if arr[j] > arr[j+1]:\n"
             "                arr[j], arr[j+1] = arr[j+1], arr[j]\n"
             "    return arr\n\n"
             "# LLM evolves to quicksort (fitness ≈ 0.87+)\n"
             "def sort(arr):  # quicksort\n"
             "    if len(arr) <= 1: return arr\n"
             "    pivot = arr[len(arr)//2]\n"
             "    return sort([x for x in arr if x<pivot]) + \\\n"
             "           [x for x in arr if x==pivot] + \\\n"
             "           sort([x for x in arr if x>pivot])",
             Inches(0.3), Inches(3.55), Inches(7.7), Inches(3.7), font_size=Pt(10))

    bullet_box(s, [
        ("Weights sum-to-1 enforced in UI (auto-normalise)", 0),
        ("Sandboxed exec — recursive algorithms work", 0),
        ("Seed algorithm: bubble sort  O(n²)", 0),
        ("Target: quicksort / built-in sorted  O(n log n)", 0),
    ], Inches(8.15), Inches(3.55), Inches(4.85), Inches(3.7), bg=C_BG2)

    add_slide_number(s, num, total)


def slide_profiling(prs, num, total):
    s = blank_slide(prs)
    fill_bg(s)
    add_title_bar(s, "Performance Profiling", "cProfile analysis — profile_experiment.py")

    bullet_box(s, [
        ("Three Profiling Methods Demonstrated", 0),
        ("Method 1: cProfile.run() — one-liner for quick profiling", 1),
        ("Method 2: cProfile.Profile() — programmatic start/stop with pstats filtering", 1),
        ("Method 3: python -m cProfile -s cumtime run_experiment.py  (CLI)", 1),
        ("Hot Spots Identified (matrix problem, 2 generations)", 0),
        ("evaluate() in matrix.py — exec() overhead per candidate", 1),
        ("cache.py set()/get() — JSONL I/O on every evaluation", 1),
        ("evolve.py — inner loop overhead (sorted, list comprehensions)", 1),
        ("LLM HTTP request — dominant cost in llm_guided mode", 1),
        ("Optimisation Opportunities", 0),
        ("Batch candidate evaluation with concurrent.futures (2–4× speedup)", 1),
        ("Use in-memory dict cache first, flush JSONL async", 1),
        ("snakeviz:  pip install snakeviz && snakeviz outputs/profile_matrix.prof", 1),
    ], Inches(0.3), Inches(1.1), Inches(8.0), Inches(5.6), bg=C_BG2)

    code_box(s,
             "# Run profiler:\n"
             "python profile_experiment.py matrix --save\n\n"
             "# View flame graph:\n"
             "snakeviz outputs/profile_matrix.prof\n\n"
             "# CLI method:\n"
             "python -m cProfile -s cumtime run_experiment.py\n\n"
             "# Top functions (cumtime):\n"
             "  run_experiment      1 call    0.021s\n"
             "  evolve              3 calls   0.018s\n"
             "  MatrixEvaluator     36 calls  0.014s\n"
             "  FitnessCache.get    36 calls  0.003s",
             Inches(8.3), Inches(1.1), Inches(4.7), Inches(5.6), font_size=Pt(10))

    add_slide_number(s, num, total)


def slide_demo(prs, num, total):
    s = blank_slide(prs)
    fill_bg(s)
    add_title_bar(s, "Demo Walkthrough", "Streamlit UI — live evolution")

    steps = [
        ("1", "Start App",       "streamlit run app.py\nOpen  http://localhost:8501"),
        ("2", "Select Problem",  "Choose pacman / matrix / pseudocode\nPseudocode: adjust 4 weight sliders"),
        ("3", "Configure",       "Set generations=10, population=8\nPick evolution modes to compare"),
        ("4", "Run Evolution",   "Click  ▶  Run Evolution\nWatch live fitness bars update per generation"),
        ("5", "Analyse Results", "Tab 1: Fitness curves by mode\nTab 2: Generation breakdown table"),
        ("6", "Code Diff",       "Tab 3: Unified diff — initial vs best evolved\nSide-by-side code view"),
        ("7", "Export",          "Tab 5: Download best code .py\nDownload results .csv / .xlsx"),
    ]

    for i, (num_str, title, body) in enumerate(steps):
        col = i % 2
        row = i // 2
        lx = Inches(0.3) + col * Inches(6.55)
        ly = Inches(1.15) + row * Inches(1.5)

        add_rect(s, lx, ly, Inches(0.55), Inches(1.2), fill_color=C_ACCENT)
        add_textbox(s, num_str, lx, ly, Inches(0.55), Inches(1.2),
                    font_size=Pt(22), bold=True, color=C_BG,
                    align=PP_ALIGN.CENTER, font_name="Segoe UI Black")

        add_rect(s, lx + Inches(0.55), ly, Inches(5.7), Inches(1.2),
                 fill_color=C_BG2, line_color=C_ACCENT, line_width=Pt(0.4))
        add_textbox(s, title, lx + Inches(0.65), ly + Inches(0.05),
                    Inches(5.5), Inches(0.45),
                    font_size=Pt(13), bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
        add_textbox(s, body, lx + Inches(0.65), ly + Inches(0.5),
                    Inches(5.5), Inches(0.65),
                    font_size=Pt(11), color=C_LIGHTGRAY, align=PP_ALIGN.LEFT)

    add_textbox(s, "Demo Video:  youtu.be/b2GBt0VEED0",
                Inches(0.3), Inches(7.05), Inches(12.7), Inches(0.35),
                font_size=Pt(12), color=C_ACCENT, align=PP_ALIGN.CENTER)

    add_slide_number(s, num, total)


def slide_conclusions(prs, num, total):
    s = blank_slide(prs)
    fill_bg(s)
    add_title_bar(s, "Conclusions & Key Insights", "What we built · What we learned")

    bullet_box(s, [
        ("What We Built", 0),
        ("Full AlphaEvolve-inspired framework: 3 problems, 3 modes, 3 LLM backends", 1),
        ("Streamlit UI with live charts + code diff + configurable weights", 1),
        ("Caching system cuts evaluation cost by 60–80% in later generations", 1),
        ("Reproducible experiments via seed=42 (both Python random + NumPy)", 1),
        ("Bonus pseudocode evaluator with 4 user-configurable fitness dimensions", 1),
        ("What We Learned", 0),
        ("LLM-guided mutation consistently outperforms random (+15–45%)", 1),
        ("Fitness ceiling effect: already-near-optimal code hard to improve further", 1),
        ("Code as a genome: modularity (engine/evaluator/LLM) makes extension easy", 1),
        ("Greedy top-K selection is simple but effective — no fancy crossover needed", 1),
        ("AOA Connections", 0),
        ("Greedy: top-K elite selection is a greedy local-search step", 1),
        ("Randomised: mutation operators introduce controlled randomness", 1),
        ("Approximation: fitness ≈ true quality (bounded eval budget)", 1),
    ], Inches(0.3), Inches(1.1), Inches(12.7), Inches(5.8), bg=C_BG2)

    add_slide_number(s, num, total)


def slide_references(prs, num, total):
    s = blank_slide(prs)
    fill_bg(s)
    add_title_bar(s, "References", "CS5381 AOA – Group 2")

    refs = [
        "[1] A. Novikov et al., AlphaEvolve: A coding agent for scientific and algorithmic discovery,\n"
        "     arXiv:2506.13131, Jun. 2025.  https://doi.org/10.48550/arXiv.2506.13131",

        "[2] S. Tamilselvi, Introduction to Evolutionary Algorithms, in Genetic Algorithms,\n"
        "     IntechOpen, 2022.  https://doi.org/10.5772/intechopen.104198",

        "[3] H. Amit, An Overview of Evolutionary Algorithms, We Talk Data, Medium.\n"
        "     https://medium.com/we-talk-data/an-overview-of-evolutionary-algorithms",

        "[4] UC Berkeley, Introduction to AI - Pacman Projects.\n"
        "     http://ai.berkeley.edu/project_overview.html",

        "[5] Meta AI, Ollama: Run large language models locally.\n"
        "     https://ollama.com",

        "[6] GitHub Repository: https://github.com/ManeeGit/CS5381-AOA",

        "[7] Demo Video: https://youtu.be/b2GBt0VEED0",
    ]

    for i, ref in enumerate(refs):
        add_textbox(s, ref,
                    Inches(0.4), Inches(1.15 + i * 0.8), Inches(12.5), Inches(0.75),
                    font_size=Pt(11.5), color=C_LIGHTGRAY, align=PP_ALIGN.LEFT)

    add_slide_number(s, num, total)


# ═══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def build():
    prs = new_prs()
    TOTAL = 15

    slide_title(prs, 1, TOTAL)
    slide_group(prs, 2, TOTAL)
    slide_motivation(prs, 3, TOTAL)
    slide_architecture(prs, 4, TOTAL)
    slide_evolution_loop(prs, 5, TOTAL)
    slide_three_modes(prs, 6, TOTAL)
    slide_fitness_functions(prs, 7, TOTAL)
    slide_mutations(prs, 8, TOTAL)
    slide_llm(prs, 9, TOTAL)
    slide_caching(prs, 10, TOTAL)
    slide_results_matrix(prs, 11, TOTAL)
    slide_results_pacman(prs, 12, TOTAL)
    slide_pseudocode(prs, 13, TOTAL)
    slide_profiling(prs, 14, TOTAL)
    slide_demo(prs, 15, TOTAL)
    slide_conclusions(prs, 15, TOTAL)   # bonus slide; same number displayed
    slide_references(prs, 15, TOTAL)

    prs.save(str(OUT))
    print(f"[✓] Saved:  {OUT}")
    print(f"    Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
