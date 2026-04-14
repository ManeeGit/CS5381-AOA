#!/usr/bin/env python3
"""
Create Group_2_round2.pptx - AOA Project Round 2 Presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Colour palette ────────────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x0D, 0x1B, 0x2A)   # dark navy
ACCENT    = RGBColor(0x00, 0xC8, 0xFF)   # electric blue
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0x10, 0x2A, 0x45)   # card blue
GREEN     = RGBColor(0x00, 0xE6, 0x76)
YELLOW    = RGBColor(0xFF, 0xD6, 0x00)
RED       = RGBColor(0xFF, 0x4D, 0x4D)
GREY_TEXT = RGBColor(0xCC, 0xD6, 0xE0)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK_LAYOUT = prs.slide_layouts[6]   # completely blank

BASE = os.path.dirname(os.path.abspath(__file__))
IMG = {
    "bhanu_fit"    : os.path.join(BASE, "student_data/Bhanu_Sankar_Ravi/matrix_fitness_comparison.png"),
    "bhanu_det"    : os.path.join(BASE, "student_data/Bhanu_Sankar_Ravi/matrix_detailed_analysis.png"),
    "maneesh_fit"  : os.path.join(BASE, "student_data/Maneesh_Malepati/matrix_fitness_comparison.png"),
    "maneesh_det"  : os.path.join(BASE, "student_data/Maneesh_Malepati/matrix_detailed_analysis.png"),
    "lakshman_fit" : os.path.join(BASE, "student_data/Lakshman_Pukhraj/matrix_fitness_comparison.png"),
    "lakshman_det" : os.path.join(BASE, "student_data/Lakshman_Pukhraj/matrix_detailed_analysis.png"),
    "johitha_fit"  : os.path.join(BASE, "student_data/Johitha_Konduru/matrix_fitness_comparison.png"),
    "johitha_det"  : os.path.join(BASE, "student_data/Johitha_Konduru/matrix_detailed_analysis.png"),
    "matrix_out"   : os.path.join(BASE, "outputs/matrix_fitness.png"),
    "pacman_out"   : os.path.join(BASE, "outputs/pacman_fitness.png"),
}


# ── Helpers ───────────────────────────────────────────────────────────────────

def bg(slide, color=BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, x, y, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def accent_bar(slide, y_inches=0.72, h_inches=0.06):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(y_inches),
                                  W, Inches(h_inches))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

def section_header(slide, title):
    """Top accent line + slide title."""
    accent_bar(slide, y_inches=0)
    txt(slide, title, 0.4, 0.12, 12, 0.6, size=28, bold=True, color=ACCENT,
        align=PP_ALIGN.LEFT)
    accent_bar(slide, y_inches=0.72, h_inches=0.04)

def bullet_list(slide, items, x, y, w, h, size=17, color=GREY_TEXT,
                bullet="▸ ", line_h=0.42):
    for i, item in enumerate(items):
        txt(slide, bullet + item, x, y + i * line_h, w, line_h,
            size=size, color=color)

def img(slide, key, x, y, w, h=None):
    path = IMG.get(key, "")
    if not os.path.exists(path):
        return
    if h:
        slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w))

def pill(slide, label, x, y, w=2.4, h=0.45, bg_color=LIGHT_BG, fg_color=ACCENT):
    box(slide, x, y, w, h, bg_color)
    txt(slide, label, x, y+0.04, w, h, size=14, bold=True,
        color=fg_color, align=PP_ALIGN.CENTER)

def divider(slide, y):
    accent_bar(slide, y_inches=y, h_inches=0.025)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – Title
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
bg(slide)

# big accent top band
box(slide, 0, 0, 13.33, 0.5, ACCENT)

# decorative right panel
box(slide, 9.8, 0.5, 3.53, 7.0, LIGHT_BG)

txt(slide, "EVOLVE", 0.5, 0.7, 9.0, 1.5, size=72, bold=True, color=WHITE)
txt(slide, "Making Code Write Itself (Sort Of)", 0.5, 2.1, 9.0, 0.7,
    size=26, color=ACCENT, bold=True)

divider(slide, 2.85)

txt(slide, "We built a system that automatically searches for better algorithms —\n"
           "using evolution, randomness, and a local LLM running on our own machines.",
    0.5, 3.0, 9.0, 1.2, size=18, color=GREY_TEXT)

# Group info box (right panel)
txt(slide, "Group 2", 10.0, 0.9, 3.0, 0.6, size=22, bold=True,
    color=ACCENT, align=PP_ALIGN.CENTER)
txt(slide, "CS5381 – Analysis of Algorithms", 9.9, 1.5, 3.3, 0.5,
    size=12, color=GREY_TEXT, align=PP_ALIGN.CENTER)

members = [
    "Bhanu Sankar Ravi",
    "Maneesh Malepati",
    "Lakshman Pukhraj",
    "Johitha Konduru",
]
for i, m in enumerate(members):
    txt(slide, m, 9.9, 2.2 + i*0.52, 3.3, 0.48,
        size=14, color=WHITE, align=PP_ALIGN.CENTER)

txt(slide, "April 13, 2026", 9.9, 6.6, 3.3, 0.5,
    size=13, color=ACCENT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – Problem Statement (1)
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
bg(slide)
section_header(slide, "The Problem We Wanted to Solve")

txt(slide, "Can a computer discover better algorithms on its own?",
    0.4, 0.9, 12, 0.6, size=22, bold=True, color=WHITE)

cols = [
    ("Hand-tuning is slow\nand painful", [
        "Engineers spend days tweaking constants by hand",
        "Our intuition is surprisingly bad at non-obvious tricks",
        "What works for one problem rarely transfers to another",
    ], RED),
    ("The search space is\npractically infinite", [
        "Even small programs have astronomically many variants",
        "You can't try everything — you need a smarter strategy",
        "Greedy hill-climbing gets stuck in local optima fast",
    ], YELLOW),
    ("LLMs guess; they\ndon't actually learn", [
        "ChatGPT-style tools produce code that looks right but isn't",
        "There's no feedback loop — each generation is independent",
        "Iterative refinement is exactly what evolution gives us",
    ], ACCENT),
]

for i, (title, bullets, color) in enumerate(cols):
    cx = 0.4 + i * 4.3
    box(slide, cx, 1.65, 4.0, 5.3, LIGHT_BG)
    bar_sh = slide.shapes.add_shape(1, Inches(cx), Inches(1.65), Inches(4.0), Inches(0.08))
    bar_sh.fill.solid(); bar_sh.fill.fore_color.rgb = color; bar_sh.line.fill.background()
    txt(slide, title, cx+0.15, 1.75, 3.7, 0.9, size=15, bold=True, color=color)
    for j, b in enumerate(bullets):
        txt(slide, "• " + b, cx+0.15, 2.7 + j*0.75, 3.7, 0.7, size=13, color=GREY_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – Problem Statement (2) – Benchmark Problems
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
bg(slide)
section_header(slide, "What We Actually Tested On")

txt(slide, "We picked two very different problems to test whether Evolve generalises",
    0.4, 0.9, 12, 0.55, size=22, bold=True, color=WHITE)

# Pacman column
box(slide, 0.4, 1.6, 6.0, 5.4, LIGHT_BG)
txt(slide, "  Pacman Agent", 0.6, 1.65, 5.7, 0.55, size=19, bold=True, color=ACCENT)
divider(slide, 2.22)
txt(slide, "Can Evolve teach Pacman to stay alive longer and score more points? "
           "The agent is pure Python — every generation mutates the decision logic.",
    0.6, 2.3, 5.7, 0.7, size=14, color=GREY_TEXT)
txt(slide, "How we score each agent:", 0.6, 3.1, 5.7, 0.4, size=14, bold=True, color=WHITE)
box(slide, 0.55, 3.5, 5.75, 0.65, BG_DARK)
txt(slide, "F = 0.6 × score  +  0.3 × survival_time  −  0.1 × steps",
    0.65, 3.53, 5.6, 0.55, size=14, bold=True, color=GREEN)
txt(slide, "We care most about score, then survival. Wasting moves is penalised.",
    0.6, 4.3, 5.7, 1.0, size=13, color=GREY_TEXT)

# Matrix column
box(slide, 6.9, 1.6, 6.0, 5.4, LIGHT_BG)
txt(slide, "  3×3 Matrix Multiply", 7.1, 1.65, 5.7, 0.55, size=19, bold=True, color=ACCENT)
divider(slide, 2.22)
txt(slide, "Can Evolve discover a faster-than-naïve way to multiply two 3×3 matrices? "
           "Inspired directly by DeepMind's AlphaTensor result.",
    7.1, 2.3, 5.7, 0.7, size=14, color=GREY_TEXT)
txt(slide, "How we score each solution:", 7.1, 3.1, 5.7, 0.4, size=14, bold=True, color=WHITE)
box(slide, 6.85, 3.5, 5.75, 0.65, BG_DARK)
txt(slide, "F = 0.7 × correctness  +  0.3 × (1 − normalized_ops)",
    6.95, 3.53, 5.6, 0.55, size=14, bold=True, color=GREEN)
txt(slide, "Correctness is non-negotiable. Op reduction is the efficiency reward.",
    7.1, 4.3, 5.7, 1.0, size=13, color=GREY_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – Proposed Solution (1)
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
bg(slide)
section_header(slide, "How Evolve Works")

txt(slide, "Six simple steps that repeat until nothing gets better",
    0.4, 0.9, 12, 0.55, size=22, bold=True, color=WHITE)

steps = [
    ("1", "Start", "Seed the population with the base code"),
    ("2", "Run it", "Execute every candidate, measure the fitness score"),
    ("3", "Keep the best", "Survive only the Top-K; drop the rest"),
    ("4", "Make new versions", "Random tweaks, template swaps, or LLM-rewrites"),
    ("5", "Repeat", "Loop for N generations — or stop early if stuck"),
    ("6", "Done", "Return the best candidate found, with full history"),
]
for i, (num, title, desc) in enumerate(steps):
    col = i % 3
    row = i // 3
    cx  = 0.4 + col * 4.3
    cy  = 1.7 + row * 2.6
    box(slide, cx, cy, 4.0, 2.2, LIGHT_BG)
    # number circle
    circ = slide.shapes.add_shape(9, Inches(cx+0.15), Inches(cy+0.15),
                                   Inches(0.6), Inches(0.6))
    circ.fill.solid(); circ.fill.fore_color.rgb = ACCENT; circ.line.fill.background()
    txt(slide, num, cx+0.15, cy+0.15, 0.6, 0.55, size=16, bold=True,
        color=BG_DARK, align=PP_ALIGN.CENTER)
    txt(slide, title, cx+0.9, cy+0.2, 3.0, 0.5, size=16, bold=True, color=WHITE)
    txt(slide, desc,  cx+0.15, cy+0.9, 3.75, 0.8, size=13, color=GREY_TEXT)

# arrow between rows
arr = slide.shapes.add_shape(1, Inches(6.0), Inches(3.75), Inches(1.3), Inches(0.08))
arr.fill.solid(); arr.fill.fore_color.rgb = ACCENT; arr.line.fill.background()
txt(slide, "↓", 6.4, 3.6, 0.5, 0.35, size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – Proposed Solution (2) – Three Evolution Modes
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
bg(slide)
section_header(slide, "Three Modes — From Dumb to Smart")

modes = [
    ("no_evolution", "Do Nothing (Baseline)", [
        "Run the base code, score it, call it a day",
        "Our honest lower-bound — how good is the starting code?",
        "Spoiler: it's already pretty good on matrix multiply",
    ], GREY_TEXT),
    ("random_mutation", "Random Mutation", [
        "Perturb numeric constants at random",
        "Swap two random lines and see what happens",
        "Inject code snippets from a hand-crafted template library",
    ], YELLOW),
    ("llm_guided", "LLM-Guided (our favourite)", [
        "Feed the code + the fitness goal to a local Ollama model",
        "qwen2.5-coder:1.5b rewrites it on your own hardware",
        "No cloud, no API key, no data leaving your machine",
    ], GREEN),
]

for i, (mode_id, title, bullets, color) in enumerate(modes):
    cx = 0.4 + i * 4.3
    box(slide, cx, 1.6, 4.0, 5.4, LIGHT_BG)
    bar = slide.shapes.add_shape(1, Inches(cx), Inches(1.6), Inches(4.0), Inches(0.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
    txt(slide, title, cx+0.15, 1.75, 3.7, 0.55, size=18, bold=True, color=color)
    box(slide, cx+0.15, 2.35, 3.7, 0.38, BG_DARK)
    txt(slide, mode_id, cx+0.2, 2.37, 3.6, 0.35, size=12, italic=True, color=ACCENT)
    for j, b in enumerate(bullets):
        txt(slide, "▸ " + b, cx+0.15, 2.85 + j*0.75, 3.7, 0.7, size=13, color=GREY_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – System Requirements (1)
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
bg(slide)
section_header(slide, "What We Used to Build It")

left_items = [
    ("Python 3.9+", "Everything runs in plain Python"),
    ("Streamlit",   "The web UI — surprisingly easy to use"),
    ("FastAPI",     "REST backend for programmatic access"),
    ("Pandas / NumPy", "Results analysis and number crunching"),
    ("Matplotlib",  "Generating all the fitness plots"),
    ("python-pptx", "Yes, this very presentation was generated by code"),
]
right_items = [
    ("Ollama",      "Runs the LLM entirely on your own machine"),
    ("qwen2.5-coder:1.5b", "Our main model — small, fast, surprisingly capable"),
    ("qwen2.5:7b",  "Heavier fallback when we need more reasoning"),
    ("UC Berkeley Pacman", "Classic CS teaching framework repurposed for eval"),
    ("SHA-256 Cache", "Avoids re-running code we've already evaluated"),
    ("YAML Config",  "One config file controls the whole experiment"),
]

txt(slide, "App & Tooling", 0.4, 0.95, 6.0, 0.45, size=16, bold=True, color=ACCENT)
txt(slide, "AI & Evaluation", 6.9, 0.95, 6.0, 0.45, size=16, bold=True, color=ACCENT)

for i, (lib, desc) in enumerate(left_items):
    cy = 1.5 + i * 0.88
    box(slide, 0.4, cy, 6.1, 0.75, LIGHT_BG)
    txt(slide, lib,  0.6, cy+0.06, 2.8, 0.55, size=15, bold=True, color=WHITE)
    txt(slide, desc, 3.5, cy+0.06, 3.0, 0.55, size=13, color=GREY_TEXT)

for i, (lib, desc) in enumerate(right_items):
    cy = 1.5 + i * 0.88
    box(slide, 6.9, cy, 6.1, 0.75, LIGHT_BG)
    txt(slide, lib,  7.1, cy+0.06, 3.0, 0.55, size=15, bold=True, color=WHITE)
    txt(slide, desc, 10.2, cy+0.06, 2.7, 0.55, size=13, color=GREY_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – System Requirements (2) – Dependencies
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
bg(slide)
section_header(slide, "Getting It Running (Under 5 Minutes)")

txt(slide, "Four steps and you're ready to go", 0.4, 0.9, 12, 0.5, size=20, bold=True, color=WHITE)

prereqs = [
    ("Python 3.9+",      "macOS / Linux / Windows — any will do",       GREEN),
    ("Ollama",           "brew install ollama  (macOS)",   ACCENT),
    ("Virtual Env",      "python3 -m venv .venv",          YELLOW),
    ("pip dependencies", "pip install -r requirements.txt", GREY_TEXT),
]

for i, (name, cmd, color) in enumerate(prereqs):
    cy = 1.55 + i * 1.3
    box(slide, 0.4, cy, 5.8, 1.1, LIGHT_BG)
    bar = slide.shapes.add_shape(1, Inches(0.4), Inches(cy), Inches(0.08), Inches(1.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
    txt(slide, name, 0.65, cy+0.1, 5.3, 0.4, size=15, bold=True, color=WHITE)
    box(slide, 0.55, cy+0.55, 5.55, 0.42, BG_DARK)
    txt(slide, cmd,  0.65, cy+0.57, 5.3, 0.38, size=13, italic=True, color=ACCENT)

# Quick start block
box(slide, 6.9, 1.55, 6.1, 5.4, LIGHT_BG)
txt(slide, "Copy-paste this and you're live", 7.1, 1.65, 5.7, 0.45, size=16, bold=True, color=ACCENT)

cmds = [
    "# 1. Clone & setup",
    "git clone <repo>  &&  cd Project",
    "python3 -m venv .venv  &&  source .venv/bin/activate",
    "pip install -r requirements.txt",
    "",
    "# 2. Start local LLM",
    "ollama serve  &",
    "ollama pull qwen2.5-coder:1.5b",
    "",
    "# 3. Run UI",
    "streamlit run app.py",
    "",
    "# 4. Run experiment",
    "python3 collect_student_data.py \\",
    "  --student 'Your Name' --config-id 1",
]
for i, cmd in enumerate(cmds):
    color = ACCENT if cmd.startswith("#") else (WHITE if cmd else WHITE)
    txt(slide, cmd, 7.1, 2.2 + i*0.28, 5.7, 0.26, size=10.5,
        color=color, italic=cmd.startswith("#"))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – Architecture Diagram (1) – High-Level
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
bg(slide)
section_header(slide, "The Big Picture — How Everything Connects")

layers = [
    ("You talk to one of these:", [
        ("Streamlit UI  (app.py)", ACCENT),
        ("FastAPI  (api.py)", ACCENT),
        ("CLI  (run_experiment.py)", ACCENT),
    ], 1.4),
    ("They all funnel into the engine:", [
        ("runner.py  –  kicks off the experiment", WHITE),
        ("evolve.py  –  the actual evolutionary loop", WHITE),
        ("mutations.py  –  decides how to mutate", WHITE),
    ], 2.75),
    ("The engine leans on these helpers:", [
        ("LLM  (src/llm/)", YELLOW),
        ("Evaluators  (src/evaluators/)", GREEN),
        ("Cache  (src/cache/)", GREY_TEXT),
    ], 4.1),
]

for layer_title, components, cy in layers:
    box(slide, 0.4, cy, 12.5, 1.2, LIGHT_BG)
    txt(slide, layer_title, 0.55, cy+0.06, 3.5, 0.35, size=13, bold=True, color=ACCENT)
    for j, (comp, color) in enumerate(components):
        cx = 3.8 + j * 3.1
        box(slide, cx, cy+0.1, 2.9, 0.95, BG_DARK)
        txt(slide, comp, cx+0.1, cy+0.3, 2.7, 0.55, size=12, color=color, align=PP_ALIGN.CENTER)
    # downward arrow
    if cy < 4.1:
        txt(slide, "↓", 6.4, cy+1.22, 0.5, 0.35, size=18, bold=True,
            color=ACCENT, align=PP_ALIGN.CENTER)

# Data flow box
box(slide, 0.4, 5.55, 12.5, 1.55, LIGHT_BG)
txt(slide, "End-to-end:", 0.55, 5.62, 2.0, 0.4, size=13, bold=True, color=ACCENT)
flow = "Config → Seed Population → Run & Score → Keep Top-K → Mutate → Next Generation → … → Best Code Found"
txt(slide, flow, 0.55, 6.05, 12.2, 0.5, size=13, color=GREY_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – Architecture Diagram (2) – Component Details
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
bg(slide)
section_header(slide, "Under the Hood — Key Files Explained")

comps = [
    ("evolve.py", "The Heart", [
        "Holds the EvolutionConfig settings",
        "Evolver.run() — one method drives everything",
        "Detects when the run is stuck and stops early",
        "Returns a clean RunResult with full history",
    ], GREEN),
    ("mutations.py", "How Code Changes", [
        "random_perturb_parameters() — nudge numbers",
        "swap_two_lines() — shuffle the logic",
        "replace_fragment() — inject known-good snippets",
        "mutate_with_meta() — ask the LLM for a rewrite",
    ], YELLOW),
    ("cache.py", "Don't Re-Run What We've Seen", [
        "Hash every code variant with SHA-256",
        "Store results in a local JSONL file",
        "TF-IDF similarity lookup for near-duplicates",
        "Survives across runs — speeds up repeat experiments",
    ], ACCENT),
    ("evaluators/", "The Judge", [
        "pacman.py — launches the game, reads the score",
        "matrix.py — counts ops, checks correctness",
        "wrappers.py — same interface for both problems",
        "Runs in a subprocess so crashes don't kill us",
    ], RED),
]

for i, (filename, title, bullets, color) in enumerate(comps):
    cx = 0.4 + i * 3.2
    box(slide, cx, 0.9, 3.0, 6.1, LIGHT_BG)
    bar = slide.shapes.add_shape(1, Inches(cx), Inches(0.9), Inches(3.0), Inches(0.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
    txt(slide, filename, cx+0.1, 1.02, 2.8, 0.42, size=13, bold=True, color=color)
    txt(slide, title,    cx+0.1, 1.48, 2.8, 0.38, size=14, bold=True, color=WHITE)
    divider(slide, 1.92)
    for j, b in enumerate(bullets):
        txt(slide, "• " + b, cx+0.1, 2.0 + j*0.8, 2.8, 0.75, size=12, color=GREY_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – Data Collection & Analysis (1) – Configurations
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
bg(slide)
section_header(slide, "Our Experiments — Who Ran What")

txt(slide, "Each of us ran a different hyperparameter setup so we could compare across configs:",
    0.4, 0.9, 12, 0.45, size=16, color=GREY_TEXT)

headers = ["Member", "Config", "Generations", "Pop Size", "Top-K", "Mut Rate"]
col_x   = [0.35, 3.1, 5.6, 7.45, 9.3, 10.9]
col_w   = [2.7,  2.4, 1.8, 1.8,  1.5, 1.8]

# Header row
box(slide, 0.35, 1.45, 12.6, 0.5, ACCENT)
for j, h in enumerate(headers):
    txt(slide, h, col_x[j]+0.05, 1.47, col_w[j], 0.42,
        size=13, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

rows = [
    ["Bhanu Sankar Ravi",  "Fast Explorer",     "8",  "6",  "2", "0.60"],
    ["Maneesh Malepati",   "Balanced",          "10", "8",  "3", "0.35"],
    ["Lakshman Pukhraj",   "Careful & Slow",    "12", "10", "4", "0.20"],
    ["Johitha Konduru",    "Elite-Focused",     "10", "12", "6", "0.30"],
]

row_colors = [LIGHT_BG, BG_DARK, LIGHT_BG, BG_DARK]
for i, row in enumerate(rows):
    cy = 2.0 + i * 0.82
    box(slide, 0.35, cy, 12.6, 0.75, row_colors[i])
    for j, cell in enumerate(row):
        txt(slide, cell, col_x[j]+0.05, cy+0.1, col_w[j], 0.55,
            size=13, color=WHITE, align=PP_ALIGN.CENTER)

txt(slide, "Matrix Multiplication  |  LLM: qwen2.5-coder:1.5b via Ollama (runs fully offline)",
    0.4, 5.5, 12.5, 0.4, size=12, italic=True, color=GREY_TEXT, align=PP_ALIGN.CENTER)

# fitness score achieved
txt(slide, "Every config hit the same peak fitness: 0.99 — the starting code was already very good",
    0.4, 5.95, 12.5, 0.45, size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

box(slide, 0.35, 6.5, 12.6, 0.7, LIGHT_BG)
txt(slide,
    "Fitness = 0.7 × correctness  +  0.3 × (1 − normalized_ops)"
    "   →   0.7 × 1.0  +  0.3 × (1 − 2/60)  ≈  0.99",
    0.5, 6.52, 12.4, 0.6, size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 – Data Collection & Analysis (2) – Fitness Comparison Plots
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
bg(slide)
section_header(slide, "Our Results — Fitness Across All Four Setups")

txt(slide, "Three lines per chart: baseline  |  random mutation  |  LLM-guided",
    0.4, 0.85, 12, 0.4, size=14, color=GREY_TEXT, align=PP_ALIGN.CENTER)

plot_data = [
    ("bhanu_fit",    "Bhanu – Fast Exploration",  0.35, 1.35),
    ("maneesh_fit",  "Maneesh – Balanced",         6.85, 1.35),
    ("lakshman_fit", "Lakshman – Conservative",    0.35, 4.2),
    ("johitha_fit",  "Johitha – Elite-Focused",    6.85, 4.2),
]

for key, label, px, py in plot_data:
    img(slide, key, px, py, 6.1, 3.0)
    txt(slide, label, px, py-0.38, 6.1, 0.35, size=12, bold=True,
        color=ACCENT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 – Data Collection & Analysis (3) – Detailed Plots
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
bg(slide)
section_header(slide, "Zooming In — Bhanu & Maneesh")

img(slide, "bhanu_det",   0.35, 1.0,  6.3, 2.95)
img(slide, "maneesh_det", 6.7,  1.0,  6.3, 2.95)

txt(slide, "Bhanu  —  Fast Explorer: high mutation (μ=0.6), small pop (6), 8 generations",
    0.35, 4.05, 6.3, 0.38, size=11, color=GREY_TEXT, align=PP_ALIGN.CENTER)
txt(slide, "Maneesh  —  Balanced: moderate mutation (μ=0.35), pop=8, 10 generations",
    6.7,  4.05, 6.3, 0.38, size=11, color=GREY_TEXT, align=PP_ALIGN.CENTER)

# Observations
box(slide, 0.35, 4.55, 12.6, 2.6, LIGHT_BG)
txt(slide, "What we noticed", 0.55, 4.62, 12, 0.4, size=15, bold=True, color=ACCENT)
obs = [
    "All three modes hit 0.99 — the starting code was essentially already optimal for this task.",
    "The baseline scores 0.99 at generation 0. That's either impressive seeding or a near-trivial problem for NumPy.",
    "Higher mutation rate (Bhanu, μ=0.6) finishes faster; lower rate (Lakshman, μ=0.2) runs more generations for the same end result.",
    "The LLM rewrites look different but score the same — it needs fitness feedback in the prompt to actually improve things.",
]
for i, o in enumerate(obs):
    txt(slide, "• " + o, 0.55, 5.1 + i*0.47, 12.3, 0.44, size=12, color=GREY_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 – Data Collection & Analysis (4) – Runtime & Mode Comparison
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
bg(slide)
section_header(slide, "Zooming In — Lakshman & Johitha + Numbers")

# Runtime table
txt(slide, "How fast did each run complete?", 0.4, 0.9, 6.2, 0.42, size=17, bold=True, color=WHITE)

rt_headers = ["Member", "Generations", "Evaluations", "Runtime (s)"]
rt_col_x   = [0.35, 2.7, 4.8, 6.5]
rt_col_w   = [2.3,  2.0, 1.9, 1.9]
box(slide, 0.35, 1.38, 8.1, 0.45, ACCENT)
for j, h in enumerate(rt_headers):
    txt(slide, h, rt_col_x[j]+0.05, 1.4, rt_col_w[j], 0.38,
        size=12, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

rt_rows = [
    ["Bhanu",   "8",  "17", "0.020"],
    ["Maneesh", "10", "21", "0.013"],
    ["Lakshman","12", "25", "0.010"],
    ["Johitha", "10", "21", "0.012"],
]
for i, row in enumerate(rt_rows):
    cy = 1.88 + i * 0.6
    box(slide, 0.35, cy, 8.1, 0.55, [LIGHT_BG, BG_DARK][i%2])
    for j, cell in enumerate(row):
        txt(slide, cell, rt_col_x[j]+0.05, cy+0.08, rt_col_w[j], 0.38,
            size=12, color=WHITE, align=PP_ALIGN.CENTER)

# Mode performance
txt(slide, "Does the mutation strategy actually matter?", 0.4, 4.4, 6.2, 0.42, size=17, bold=True, color=WHITE)
mode_rows = [
    ("no_evolution",   "0.9900", "0.00%", GREY_TEXT),
    ("random_mutation","0.9900", "0.00%", YELLOW),
    ("llm_guided",     "0.9900", "0.00%", GREEN),
]
ml_headers = ["Mode", "Best Fitness", "Improvement"]
ml_col_x   = [0.35, 3.5, 6.0]
ml_col_w   = [3.1,  2.4, 2.2]
box(slide, 0.35, 4.88, 8.1, 0.42, ACCENT)
for j, h in enumerate(ml_headers):
    txt(slide, h, ml_col_x[j]+0.05, 4.9, ml_col_w[j], 0.35,
        size=12, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
for i, (mode, fit, imp, color) in enumerate(mode_rows):
    cy = 5.35 + i*0.58
    box(slide, 0.35, cy, 8.1, 0.52, [LIGHT_BG, BG_DARK][i%2])
    for j, (cell, cw, cx) in enumerate(zip([mode, fit, imp], ml_col_w, ml_col_x)):
        txt(slide, cell, cx+0.05, cy+0.08, cw, 0.38,
            size=12, color=color, align=PP_ALIGN.CENTER)

# Lakshman & Johitha detailed plots on right
img(slide, "lakshman_det", 8.6, 0.85, 4.7, 2.22)
txt(slide, "Lakshman  —  Careful & Slow (12 gens, pop=10)", 8.6, 3.1, 4.7, 0.32, size=10, color=GREY_TEXT, align=PP_ALIGN.CENTER)
img(slide, "johitha_det", 8.6, 3.55, 4.7, 2.22)
txt(slide, "Johitha  —  Elite-Focused (keeps top 6 of 12)", 8.6, 5.8, 4.7, 0.32, size=10, color=GREY_TEXT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 – Pros and Cons
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
bg(slide)
section_header(slide, "What Worked, What Didn't")

pros = [
    "Runs 100% offline — your code never leaves your machine",
    "Adding a new problem takes about 50 lines of code",
    "Caching means we don't waste time re-running what we know",
    "Three modes let us isolate exactly what's helping",
    "One YAML file controls everything — no hardcoded magic numbers",
    "Works on games AND math problems with the same engine",
    "The UI is genuinely usable by non-programmers",
]
cons = [
    "We picked a benchmark the starting code was already great at — not much room to improve",
    "The LLM doesn't see the fitness score when it rewrites — it's essentially flying blind",
    "You still have to install Ollama yourself; we haven't bundled it",
    "Pacman evaluation drags in the full UC Berkeley game framework",
    "High mutation rate makes everyone look the same after a few generations",
    "Fitness is one number — we can't trade off speed vs. correctness yet",
]

# PRO column
box(slide, 0.35, 1.0, 5.85, 6.1, LIGHT_BG)
bar = slide.shapes.add_shape(1, Inches(0.35), Inches(1.0), Inches(5.85), Inches(0.1))
bar.fill.solid(); bar.fill.fore_color.rgb = GREEN; bar.line.fill.background()
txt(slide, "  What we're proud of", 0.55, 1.12, 5.5, 0.45, size=18, bold=True, color=GREEN)
for i, p in enumerate(pros):
    txt(slide, "  " + p, 0.55, 1.65 + i*0.75, 5.55, 0.68, size=13, color=GREY_TEXT)

# CON column
box(slide, 7.15, 1.0, 5.85, 6.1, LIGHT_BG)
bar2 = slide.shapes.add_shape(1, Inches(7.15), Inches(1.0), Inches(5.85), Inches(0.1))
bar2.fill.solid(); bar2.fill.fore_color.rgb = RED; bar2.line.fill.background()
txt(slide, "  Honest limitations", 7.35, 1.12, 5.5, 0.45, size=18, bold=True, color=RED)
for i, c in enumerate(cons):
    txt(slide, "  " + c, 7.35, 1.65 + i*0.75, 5.5, 0.68, size=13, color=GREY_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 – Issues & Solutions
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
bg(slide)
section_header(slide, "Things That Broke (and How We Fixed Them)")

issues = [
    (
        "Ollama wasn't running when experiments kicked off",
        "The LLM mode silently fell back to random mutation — we got no AI rewrites at all",
        "Added a health-check before each run; now it warns you loudly and logs the fallback",
        YELLOW,
    ),
    (
        "The matrix problem was already basically solved",
        "NumPy's einsum uses only 2 operations — there was almost nothing left to optimise",
        "We kept it — fitness of 0.99 is actually a useful result; it confirms the baseline is near-optimal",
        ACCENT,
    ),
    (
        "pandoc couldn't find our PNG files when generating DOCX reports",
        "It resolved image paths relative to the wrong directory and produced broken documents",
        "Ran pandoc from inside each student folder using pushd/popd — paths resolved correctly",
        GREEN,
    ),
    (
        "App crashed when launched from a different working directory",
        "All file paths were relative to CWD — worked in dev, broke in any other launch context",
        "Replaced every relative path with Path(__file__).parent anchored paths throughout",
        GREY_TEXT,
    ),
]

for i, (issue, impact, solution, color) in enumerate(issues):
    cy = 1.0 + i * 1.55
    box(slide, 0.35, cy, 12.6, 1.4, LIGHT_BG)
    bar = slide.shapes.add_shape(1, Inches(0.35), Inches(cy), Inches(0.08), Inches(1.4))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
    txt(slide, f"Problem: {issue}", 0.6, cy+0.05, 12.0, 0.38, size=14, bold=True, color=WHITE)
    txt(slide, f"What happened:  {impact}",   0.6, cy+0.48, 12.0, 0.38, size=12, color=RED)
    txt(slide, f"How we fixed it:  {solution}", 0.6, cy+0.9,  12.0, 0.42, size=12, color=GREEN)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 – Conclusion
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
bg(slide)
section_header(slide, "So, What Did We Actually Learn?")

txt(slide,
    "Evolve showed us that you don't need a massive compute cluster or a cloud bill\n"
    "to do meaningful automated algorithm search. A small LLM on your laptop +\n"
    "a tight feedback loop gets you surprisingly far.",
    0.4, 0.95, 12.5, 1.0, size=18, color=WHITE)

takeaways = [
    ("The loop works", "Three modes, two problem domains, four configs — results are consistent and reproducible"),
    ("Local LLMs are real now", "qwen2.5-coder:1.5b on Ollama — runs on your laptop, no API key, no latency tax"),
    ("Caching matters more than you'd think", "Deduplication cut repeated evaluation time to near-zero across runs"),
    ("Sometimes the answer is 'it's already good'", "F=0.99 on matrix multiply is a result — it tells us the baseline is near-optimal"),
    ("Modularity pays off", "Swapping in a new problem domain took under an hour because the interfaces are clean"),
]

for i, (title, desc) in enumerate(takeaways):
    cy = 2.15 + i * 0.97
    box(slide, 0.35, cy, 12.6, 0.85, LIGHT_BG)
    circ = slide.shapes.add_shape(9, Inches(0.42), Inches(cy+0.12),
                                   Inches(0.55), Inches(0.55))
    circ.fill.solid(); circ.fill.fore_color.rgb = ACCENT; circ.line.fill.background()
    txt(slide, "", 0.42, cy+0.12, 0.55, 0.5, size=14, bold=True,
        color=BG_DARK, align=PP_ALIGN.CENTER)
    txt(slide, title, 1.1, cy+0.08, 4.0, 0.38, size=14, bold=True, color=WHITE)
    txt(slide, desc,  5.2, cy+0.1,  7.6, 0.62, size=13, color=GREY_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 – Future Directions
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
bg(slide)
section_header(slide, "Where We'd Take This Next")

directions = [
    ("Pick harder problems",
     "Sorting, graph search, NP-hard puzzles — places where the starting code isn't already 99% optimal",
     ACCENT),
    ("Optimise for multiple things at once",
     "Right now fitness is one number. Real tradeoffs — speed vs. memory vs. accuracy — need a Pareto-front",
     GREEN),
    ("Tell the LLM its score",
     "The model currently doesn't know how well the previous version did. That's the obvious missing piece",
     YELLOW),
    ("Add crossover — not just mutation",
     "Splice two high-fitness parents together. Classic genetic algorithm, and we haven't tried it yet",
     ACCENT),
    ("Run evaluations in parallel",
     "Each candidate is independent — there's no reason to wait for one before starting the next",
     GREEN),
    ("Evolve the hyperparameters too",
     "Meta-evolution: let the system discover its own mutation rate and population size",
     YELLOW),
]

for i, (title, desc, color) in enumerate(directions):
    col = i % 2
    row = i // 2
    cx = 0.35 + col * 6.65
    cy = 1.05 + row * 2.15
    box(slide, cx, cy, 6.3, 2.0, LIGHT_BG)
    bar = slide.shapes.add_shape(1, Inches(cx), Inches(cy), Inches(6.3), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
    txt(slide, title, cx+0.15, cy+0.15, 6.0, 0.42, size=15, bold=True, color=WHITE)
    txt(slide, desc,  cx+0.15, cy+0.65, 6.0, 1.1,  size=12, color=GREY_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 – References
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
bg(slide)
section_header(slide, "References")

refs = [
    '[1] Fawzi, A. et al. (2022). "Discovering faster matrix multiplication algorithms with reinforcement\n'
    '     learning." Nature 610, 47–53. https://doi.org/10.1038/s41586-022-05172-4',

    '[2] Google DeepMind (2025). "AlphaEvolve: A Gemini-based coding agent for designing advanced\n'
    '     algorithms." DeepMind Blog. https://deepmind.google/discover/blog/alphaevolve/',

    '[3] Chen, M. et al. (2021). "Evaluating Large Language Models Trained on Code." OpenAI.\n'
    '     arXiv:2107.03374.',

    '[4] UC Berkeley AI Project (2023). "The Pacman Projects."\n'
    '     https://ai.berkeley.edu/project_overview.html',

    '[5] Ollama (2024). "Run LLMs locally." https://ollama.com',

    '[6] Real, E. et al. (2019). "Regularized Evolution for Image Classifier Architecture Search."\n'
    '     AAAI 2019. arXiv:1802.01548.',

    '[7] Back, T. (1996). "Evolutionary Algorithms in Theory and Practice." Oxford University Press.',
]

for i, ref in enumerate(refs):
    txt(slide, ref, 0.4, 1.0 + i * 0.85, 12.5, 0.82, size=12, color=GREY_TEXT)


# ── Save ──────────────────────────────────────────────────────────────────────
OUT = os.path.join(os.path.dirname(BASE), "Group_2_round2.pptx")
prs.save(OUT)
print(f" Presentation saved: {OUT}")
print(f"  Slides: {len(prs.slides)}")
